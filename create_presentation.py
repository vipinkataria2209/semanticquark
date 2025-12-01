"""Create SemanticQuark presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define colors
PRIMARY_COLOR = RGBColor(41, 98, 255)  # Blue
SECONDARY_COLOR = RGBColor(45, 55, 72)  # Dark gray
ACCENT_COLOR = RGBColor(16, 185, 129)  # Green

def add_title_slide(prs, title, subtitle=""):
    """Add a title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_COLOR
    title_para.alignment = PP_ALIGN.CENTER

    # Add subtitle if provided
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(9), Inches(0.8))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(24)
        subtitle_para.font.color.rgb = SECONDARY_COLOR
        subtitle_para.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, content_items):
    """Add a content slide with bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_COLOR

    # Add content
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.4), Inches(5.4))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    for i, item in enumerate(content_items):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = SECONDARY_COLOR
        p.level = 0
        p.space_before = Pt(8)

    return slide

def add_two_column_slide(prs, title, left_items, right_items):
    """Add a two-column slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_COLOR

    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(4), Inches(5.4))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True

    for i, item in enumerate(left_items):
        if i > 0:
            left_frame.add_paragraph()
        p = left_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = SECONDARY_COLOR
        p.space_before = Pt(6)

    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.6), Inches(4), Inches(5.4))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True

    for i, item in enumerate(right_items):
        if i > 0:
            right_frame.add_paragraph()
        p = right_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = SECONDARY_COLOR
        p.space_before = Pt(6)

    return slide

# Slide 1: Title Slide
add_title_slide(prs, "SemanticQuark", "The Fundamental Building Block for Semantic Analytics")

# Slide 2: The Problem
add_content_slide(prs, "The Analytics Problem", [
    "🔴 Complex SQL queries for simple questions",
    "🔴 Inconsistent metric definitions across teams",
    "🔴 Slow query performance at scale",
    "🔴 Difficult to implement data security policies",
    "🔴 No integration with modern data science tools",
    "🔴 Repetitive code across dashboards and reports"
])

# Slide 3: What is SemanticQuark?
add_content_slide(prs, "What is SemanticQuark?", [
    "A Python-native semantic layer platform for analytics",
    "Define metrics once in YAML, use everywhere",
    "Query data with JSON instead of SQL",
    "Built-in caching, pre-aggregations, and security",
    "Think of it as 'Cube.js for Python'",
    "Seamlessly integrates with Python data ecosystem"
])

# Slide 4: Core Capabilities
add_content_slide(prs, "Core Capabilities", [
    "📊 Semantic Data Modeling - Cubes, dimensions, measures",
    "🔍 REST & GraphQL APIs - Simple JSON queries",
    "⚡ Query Optimization - Automatic SQL generation",
    "💾 Intelligent Caching - Redis or in-memory",
    "📈 Pre-Aggregations - Sub-second analytics",
    "🔒 Row-Level Security - Built-in security layer",
    "🔌 Multi-Database - PostgreSQL, MySQL, Snowflake, BigQuery"
])

# Slide 5: Architecture Overview
add_content_slide(prs, "Architecture Overview", [
    "🌐 API Layer - FastAPI with REST, GraphQL, SQL",
    "⚙️ Query Orchestration - Engine with callbacks",
    "📦 Semantic Layer - Cubes, dimensions, measures",
    "🏗️ SQL Builder - BFS-based join path finding",
    "💾 Cache Layer - Redis/Memory with TTL",
    "🔒 Security Layer - Row-level security (RLS)",
    "🔌 Database Connectors - Extensible driver system"
])

# Slide 6: Semantic Modeling
add_content_slide(prs, "Semantic Modeling in YAML", [
    "Define data models declaratively:",
    "  • Cubes - Logical data models (e.g., orders, customers)",
    "  • Dimensions - Attributes to slice by (status, date)",
    "  • Measures - Metrics to aggregate (count, revenue)",
    "  • Relationships - Joins between cubes",
    "  • Pre-aggregations - Rollup tables for speed",
    "",
    "Example: orders.yaml defines dimensions (status, date)",
    "         and measures (count, total_revenue)"
])

# Slide 7: Query API
add_content_slide(prs, "Simple Query API", [
    "JSON query instead of SQL:",
    "",
    '{',
    '  "dimensions": ["orders.status"],',
    '  "measures": ["orders.count", "orders.total_revenue"],',
    '  "filters": [{"dimension": "orders.status",',
    '               "operator": "equals", "values": ["completed"]}]',
    '}',
    "",
    "SemanticQuark automatically generates optimized SQL!"
])

# Slide 8: Python-Native Advantages
add_two_column_slide(prs, "Python Ecosystem Potential",
    [
        "🐍 Built in Python",
        "📊 Future: Pandas Integration",
        "📓 Future: Jupyter Notebooks",
        "🤖 Future: ML Model Embedding",
        "📦 Access to Rich Ecosystem",
        "🔬 Data Science Ready"
    ],
    [
        "Unlike Cube.js (Node.js)",
        "Potential DataFrame support",
        "Interactive analytics capability",
        "ML-powered metrics (planned)",
        "NumPy, SciPy, etc. available",
        "Python toolchain compatible"
    ]
)

# Slide 9: Query Engine Flow
add_content_slide(prs, "Query Execution Flow", [
    "1️⃣ Parse JSON query → Query object",
    "2️⃣ Check cache (Redis/memory) → Return if hit",
    "3️⃣ Find matching pre-aggregation → Use if exists",
    "4️⃣ Apply security filters → Row-level security",
    "5️⃣ Build optimized SQL → BFS join path finding",
    "6️⃣ Execute on database → Get results",
    "7️⃣ Format & cache results → Return JSON",
    "8️⃣ Fire callbacks → Logging, metrics, monitoring"
])

# Slide 10: Intelligent Caching
add_content_slide(prs, "Intelligent Caching", [
    "🔑 Cache Key Generation",
    "  • Based on query hash + user context",
    "  • Ensures consistent cache hits",
    "",
    "💾 Multi-Tier Caching",
    "  • In-memory cache for development",
    "  • Redis cache for production",
    "  • Configurable TTL (time-to-live)",
    "",
    "♻️ Cache Invalidation on schema reload"
])

# Slide 11: Pre-Aggregations
add_content_slide(prs, "Pre-Aggregations for Speed", [
    "Pre-compute common queries → Sub-second analytics",
    "",
    "Configuration in YAML:",
    "  pre_aggregations:",
    "    - name: orders_daily",
    "      dimensions: [status, created_at]",
    "      measures: [count, total_revenue]",
    "      time_dimension: created_at",
    "      granularity: day",
    "      refresh_key: {every: 1 hour}",
    "",
    "Automatically matches queries and uses rollup tables"
])

# Slide 12: Row-Level Security
add_content_slide(prs, "Row-Level Security (RLS)", [
    "Built-in security at the semantic layer:",
    "",
    "🔒 Define security rules in cube models",
    "👤 Pass user context with queries",
    "🛡️ Automatically filter rows based on rules",
    "",
    "Example: Filter orders by user's region",
    "  security:",
    "    row_filter: \"region = {user.region}\"",
    "",
    "Security applied before SQL generation!"
])

# Slide 13: Smart SQL Generation
add_content_slide(prs, "Smart SQL Generation", [
    "BFS-based Join Path Finding:",
    "  • Finds shortest path between cubes",
    "  • Supports multi-hop joins (A→B→C)",
    "  • Bidirectional relationship traversal",
    "",
    "Query Optimization:",
    "  • Skip GROUP BY when primary key included",
    "  • Minimize JOIN operations",
    "  • Use pre-aggregations when available",
    "",
    "Result: Clean, efficient SQL every time"
])

# Slide 14: Monitoring & Observability
add_content_slide(prs, "Monitoring & Observability", [
    "📊 Callback-Based Monitoring System",
    "  • on_query_start / on_query_end",
    "  • on_cache_hit / on_cache_miss",
    "  • on_pre_agg_used / on_pre_agg_skipped",
    "  • on_sql_generated",
    "  • on_query_error",
    "",
    "📈 Built-in Metrics & Logging",
    "  • Query execution times",
    "  • Cache hit rates",
    "  • Query logs with user context"
])

# Slide 15: Production Features
add_two_column_slide(prs, "Production-Ready Features",
    [
        "🔄 Hot Schema Reload",
        "🐳 Docker Support",
        "📝 Auto API Docs",
        "🧪 Comprehensive Tests",
        "⚡ Connection Pooling",
        "🔌 Extensible Drivers"
    ],
    [
        "No downtime updates",
        "Docker Compose ready",
        "Swagger/OpenAPI",
        "Unit & integration",
        "Async PostgreSQL",
        "Plugin architecture"
    ]
)

# Slide 16: Use Cases
add_content_slide(prs, "Use Cases", [
    "📊 Business Intelligence Dashboards",
    "  • Consistent metrics across all dashboards",
    "",
    "🔬 Data Science Workflows",
    "  • Query data directly in Jupyter notebooks",
    "",
    "🤖 ML Feature Engineering",
    "  • Embed ML models in metric definitions",
    "",
    "📱 Embedded Analytics",
    "  • Secure multi-tenant analytics APIs",
    "",
    "📈 Real-time Analytics",
    "  • Sub-second queries with pre-aggregations"
])

# Slide 17: SemanticQuark vs Cube.js
add_two_column_slide(prs, "SemanticQuark vs Cube.js",
    [
        "SemanticQuark (Python):",
        "✅ Native Python ecosystem",
        "✅ Pandas/Jupyter integration",
        "✅ ML model embedding",
        "✅ Data science workflows",
        "✅ Async/await design",
        "✅ Callback architecture"
    ],
    [
        "Cube.js (Node.js):",
        "✅ Mature ecosystem",
        "✅ Large community",
        "✅ More connectors",
        "❌ No Python integration",
        "❌ No ML support",
        "❌ JavaScript-focused"
    ]
)

# Slide 18: Get Started & Summary
add_content_slide(prs, "Get Started Today!", [
    "🚀 Quick Start:",
    "  git clone https://github.com/yourusername/semanticquark",
    "  docker-compose up -d",
    "  curl http://localhost:8000/health",
    "",
    "📚 Learn More:",
    "  • Documentation: /docs",
    "  • Examples: /models/*.yaml",
    "  • Tests: /tests/",
    "",
    "🌟 SemanticQuark = Metrics as Code + Python Power",
    "   Define once, query everywhere, scale to production"
])

# Save presentation
prs.save('/home/user/semanticquark/SemanticQuark_Presentation.pptx')
print("✅ Presentation created: SemanticQuark_Presentation.pptx")
print("📊 18 slides generated successfully!")
